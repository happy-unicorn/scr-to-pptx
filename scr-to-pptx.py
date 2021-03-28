import time
import numpy as np
import pyautogui as pag
import mss
import cv2
import pathlib
import argparse
import pptx
from pptx.util import Inches


def main(args):
    screen_width, screen_height = pag.size()
    output_path = pathlib.Path(args.path)
    output_path.mkdir(parents=True, exist_ok=True)

    top = left = Inches(0)
    prs = pptx.Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    time.sleep(10)

    with mss.mss() as scr:
        for i in range(args.number_of_scr):
            image = np.asarray(scr.grab({'top': 0, 'left': 0, 'width': screen_width, 'height': screen_height}))

            scr_output_path = str(output_path / f'{i}.jpg')
            cv2.imwrite(scr_output_path, image)

            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.add_picture(scr_output_path, left, top, prs.slide_width, prs.slide_height)

            pag.press('right')
            time.sleep(1)

    prs.save(output_path / f'{args.filename}.pptx')


if __name__ == '__main__':
    parser = argparse.ArgumentParser()

    parser.add_argument('-n', '--number_of_scr', required=True, type=int)
    parser.add_argument('-p', '--path', required=True, type=pathlib.Path)
    parser.add_argument('-f', '--filename', required=True, type=str)

    args = parser.parse_args()

    main(args)
